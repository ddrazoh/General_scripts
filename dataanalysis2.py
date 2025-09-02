import pandas as pd
from Bio import pairwise2
import scipy.stats as stats
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas


# Step 1: Read the data from the CSV file
data = pd.read_csv('~/Desktop/tf_trial/analysis/testdata.csv')

# Extract the sample IDs, Sanger sequences, and MinION sequences from the dataset
sample_ids = data['Sample_id']
sanger_sequences = data['Sanger']
minion_sequences = data['MinION']

# Step 2: Calculate Error Rates using pairwise alignment
minion_error_rates = []

minion_error_rates = []
pairwise_identities = []

for i, (sanger_seq, minion_seq) in enumerate(zip(sanger_sequences, minion_sequences)):
    alignment = pairwise2.align.globalms(sanger_seq, minion_seq, 2, -1, -2, -2)  # Perform pairwise alignment
    aligned_sanger_seq, aligned_minion_seq, _, _, _ = alignment[0]

    matches = sum(a == b for a, b in zip(aligned_sanger_seq, aligned_minion_seq))
    seq_length = len(aligned_sanger_seq)

    pairwise_identity = (matches / seq_length) * 100

    mismatches = sum(a != b for a, b in zip(aligned_sanger_seq, aligned_minion_seq))
    indels = aligned_minion_seq.count('-') + aligned_sanger_seq.count('-')
    total_errors = mismatches + indels

    error_rate = (total_errors / seq_length) * 100

    minion_error_rates.append(error_rate)
    pairwise_identities.append(pairwise_identity)

for i, (sanger_seq, minion_seq) in enumerate(zip(sanger_sequences, minion_sequences)):
    alignment = pairwise2.align.globalms(sanger_seq, minion_seq, 2, -1, -2, -2)  # Perform pairwise alignment
    aligned_sanger_seq, aligned_minion_seq, _, _, _ = alignment[0]

    mismatches = 0
    indels = 0
    seq_length = len(aligned_sanger_seq)

    for j in range(seq_length):
        if aligned_sanger_seq[j] != aligned_minion_seq[j]:
            if aligned_sanger_seq[j] == '-' or aligned_minion_seq[j] == '-':
                indels += 1
            else:
                mismatches += 1

    error_rate = (mismatches + indels) / seq_length * 100

    minion_error_rates.append(error_rate)

    # Print the error rate for each sample
    print(f"Sample ID: {sample_ids[i]}")
    print(f"MinION Error Rate: {error_rate:.2f}%")
    print()

# Calculate the overall mean error rate and range
mean_error_rate = sum(minion_error_rates) / len(minion_error_rates)
error_rate_range = max(minion_error_rates) - min(minion_error_rates)

# Print the overall mean error rate and range
print("Overall Mean MinION Error Rate:", f"{mean_error_rate:.2f}%")
print("MinION Error Rate Range:", f"{error_rate_range:.2f}%")

# Step 3: Statistical Analysis
t_statistic, p_value = stats.ttest_1samp(minion_error_rates, mean_error_rate)

# Print the statistical results
print("T-test p-value:", p_value)

# Generate tables
table_data = {'Sample ID': sample_ids, 'MinION Error Rate (%)': minion_error_rates}
error_rate_table = pd.DataFrame(table_data)
error_rate_table.to_csv('error_rate_table.csv', index=False)
print("\nError Rate Table:")
print(error_rate_table)


# Generate contingency tables
contingency_tables = []
for i, sample_id in enumerate(sample_ids):
    contingency_table = pd.DataFrame({'count': [minion_error_rates[i] > mean_error_rate]}, index=[f'Sample {sample_id}'])
    contingency_tables.append(contingency_table)

# Print the contingency tables
print("\nContingency Tables:")
for i, contingency_table in enumerate(contingency_tables):
    print(f"Sample ID: {sample_ids[i]}")
    print(contingency_table)
    print()
    contingency_table.to_csv(f'contingency_table_{sample_ids[i]}.csv')
    



# Generate plots
plt.figure(figsize=(12, 8))

# Histogram
plt.subplot(2, 2, 1)
plt.hist(minion_error_rates, bins=10, edgecolor='black')
plt.xlabel('MinION Error Rate (%)')
plt.ylabel('Frequency')
plt.title('Distribution of MinION Error Rates')

# Boxplot
plt.subplot(2, 2, 2)
plt.boxplot(minion_error_rates)
plt.ylabel('MinION Error Rate (%)')
plt.title('Boxplot of MinION Error Rates')

# Violin plot
plt.subplot(2, 2, 3)
plt.violinplot(minion_error_rates)
plt.ylabel('MinION Error Rate (%)')
plt.title('Violin Plot of MinION Error Rates')

# Bar plot
plt.subplot(2, 2, 4)
plt.bar(sample_ids, minion_error_rates)
plt.xlabel('Sample ID')
plt.ylabel('MinION Error Rate (%)')
plt.title('MinION Error Rates per Sample')

# Adjust the spacing between subplots
plt.tight_layout()

# Show the plots
plt.show()


# Generate plots
fig, axs = plt.subplots(2, 2, figsize=(12, 8))

# Histogram
axs[0, 0].hist(minion_error_rates, bins=10, edgecolor='black')
axs[0, 0].set_xlabel('MinION Error Rate (%)')
axs[0, 0].set_ylabel('Frequency')
axs[0, 0].set_title('Distribution of MinION Error Rates')

# Boxplot
axs[0, 1].boxplot(minion_error_rates)
axs[0, 1].set_ylabel('MinION Error Rate (%)')
axs[0, 1].set_title('Boxplot of MinION Error Rates')

# Violin plot
axs[1, 0].violinplot(minion_error_rates)
axs[1, 0].set_ylabel('MinION Error Rate (%)')
axs[1, 0].set_title('Violin Plot of MinION Error Rates')

# Bar plot
axs[1, 1].bar(sample_ids, minion_error_rates)
axs[1, 1].set_xlabel('Sample ID')
axs[1, 1].set_ylabel('MinION Error Rate (%)')
axs[1, 1].set_title('MinION Error Rates per Sample')

# Adjust the spacing between subplots
plt.tight_layout()

# Save the plots as PowerPoint (.ppt) file
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Use layout for blank slide

for ax in axs.flat:
    fig = ax.get_figure()
    canvas = FigureCanvas(fig)
    image_stream = BytesIO()
    canvas.print_png(image_stream)
    image_stream.seek(0)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(image_stream, Inches(0.5), Inches(0.5), width=Inches(8), height=Inches(6))

output_pptx_file = 'plots.pptx'
prs.save(output_pptx_file)
print(f"Plots saved as {output_pptx_file}")
